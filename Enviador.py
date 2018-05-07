#!/usr/bin/python

from pyforms.utils.settings_manager import conf;
import settings
conf+=settings

import configparser
import csv
import smtplib


from pptx import Presentation

import time
import subprocess
import os.path

from email.mime.text import MIMEText
from email.header import Header
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders

from PyQt5 import QtWidgets
import pyforms
from   pyforms          import BaseWidget
from   pyforms.controls import ControlText
from   pyforms.controls import ControlButton
from   pyforms.controls import ControlFile
from   pyforms.controls import ControlEmptyWidget

from ConfigSMTPWindow import ConfigSMTPWindow
from ConfigEmailWindow import ConfigEmailWindow
from ListaEnvios import ListaEnvios
from Envio import Envio


def preparaMail(nomeEvento, csv_row, tipo, templateParticipacao):

    caminho = os.path.dirname(templateParticipacao)
    caminhoOut = caminho + '/gerados/'

    if not os.path.exists(caminhoOut):
        os.makedirs(caminhoOut)

    odp = Presentation(templateParticipacao)

    for slide in odp.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        novo_run = p.add_run()

                        if (run.text == '$nome'):
                            novo_run.text = csv_row[0]
                        elif (run.text == '$evento'):
                            novo_run.text = nomeEvento
                        else:
                            novo_run.text = run.text

                        novo_run.font.name = run.font.name
                        novo_run.font.size = run.font.size
                        novo_run.font.bold = run.font.bold
                        novo_run.font.italic = run.font.italic

                        r = run._r
                        r.getparent().remove(r)


    arquivoNome = 'Certificado de ' + tipo + ' - '+ csv_row[0]
    arquivoPptx = caminhoOut + arquivoNome + '.pptx'
    arquivoPdf = caminhoOut + arquivoNome + '.pdf'

    odp.save(arquivoPptx)

    try:
        subprocess.call(['/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'pdf', "--outdir", caminhoOut, arquivoPptx])
        while (not os.path.isfile(arquivoPdf)):
            time.sleep(1)
        os.remove(arquivoPptx)

        return {'nome':csv_row[0], 'email': csv_row[1], 'status':'OK', 'tipo': tipo, 'arquivo':arquivoPdf}
    except subprocess.CalledProcessError as e:
        print('CalledProcessError', e)
        return {'nome': csv_row[0], 'email': csv_row[1], 'status': 'falha', 'tipo': tipo, 'arquivo': ''}



class Enviador(BaseWidget):

    def __init__(self):
        super(Enviador, self).__init__('Simple example 1')

        self._smtp = None


    def geraCertificados(self, painelEnvios, tipo, nomeEvento, templateCertificados, csvEnvio):
        with open(csvEnvio, newline='') as csvfile:
            csv_reader = csv.reader(csvfile, delimiter=',', quotechar='|')
            for csv_row in csv_reader:
                certificado = preparaMail(nomeEvento, csv_row, tipo, templateCertificados)

                envio = Envio(certificado['nome'], certificado['email'], certificado['status'], certificado['tipo'], certificado['arquivo'])
                painelEnvios.value.addEnvio(envio)

                QtWidgets.QApplication.processEvents()

            


if __name__ == "__main__":
    pyforms.start_app(Enviador, geometry=(300, 150, 400, 500))